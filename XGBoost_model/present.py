# IMPORTANT: Before running this script, you must install the python-pptx library.
# Open your terminal or command prompt and run:
# pip install python-pptx

from pptx import Presentation

def create_presentation():
    # Create a new presentation object
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    title_slide_layout = prs.slide_layouts[0] # 0 is the layout for a Title Slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "AI Enabled Blood Bank Management System"
    subtitle.text = ("Mid-Semester BTP Progress Report\n"
                     "Group 1: Biswa Prakash Dalai, Dhiraj Yadav, Subhasish Sahoo\n"
                     "Supervisor: Dr. Anjali Mohapatra\n"
                     "IIIT Bhubaneswar")

    # --- Helper function to create bullet point slides easily ---
    def add_bullet_slide(title_text, bullet_points):
        bullet_slide_layout = prs.slide_layouts[1] # 1 is the layout for Title and Content
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = title_text
        tf = body_shape.text_frame
        
        # Add the first point
        tf.text = bullet_points[0]
        
        # Add the remaining points
        for point in bullet_points[1:]:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0

    # --- Slide 2: Problem Statement ---
    add_bullet_slide(
        "The Flaw in Reactive Inventory Management",
        [
            "Biological Constraint: Blood components have strict 5-to-42 day viability periods.",
            "The Balancing Act: Overstocking leads to high biological waste; understocking risks patient lives.",
            "The Current Flaw: Existing systems are strictly reactive. They track inventory but cannot anticipate demand spikes due to seasonal diseases or holidays."
        ]
    )

    # --- Slide 3: Our Solution ---
    add_bullet_slide(
        "A Proactive, AI-Driven Architecture",
        [
            "Transition from threshold-based alerts to predictive analytics.",
            "Objective 1: Engineer a mathematical pipeline for realistic synthetic data generation.",
            "Objective 2: Train and evaluate advanced time-series models (XGBoost & Prophet) across 8 blood groups.",
            "Objective 3: Deploy the optimal ML model as a standalone, testable FastAPI microservice."
        ]
    )

    # --- Slide 4: Architecture ---
    add_bullet_slide(
        "Decoupled Microservices Design Pattern",
        [
            "Primary Web App (MERN): React, Node.js, Express, MongoDB for role-based donor/inventory CRUD operations.",
            "Intelligence Engine (FastAPI): Isolated Python microservice to prevent heavy ML calculations from blocking web server traffic.",
            "",
            "[Please insert your Architecture Diagram here in Google Slides]"
        ]
    )

    # --- Slide 5: Methodology ---
    add_bullet_slide(
        "Handling Non-Linear Time-Series Data",
        [
            "Data Generation: Simulated 3 years of daily logs incorporating baseline constraints, linear trends, and noise.",
            "Advanced Feature Engineering (For XGBoost):",
            " - Cyclical Temporal Encoding: Mapped days/months onto Sine/Cosine continuous circles.",
            " - Lag Features: Captured historical memory (e.g., lag_7).",
            " - Rolling Volatility: 7-day rolling means to track trajectory."
        ]
    )

    # --- Slide 6: Evaluation ---
    add_bullet_slide(
        "Comparative Analysis: XGBoost vs. Prophet",
        [
            "Evaluated on an unseen 6-month test set.",
            "Prophet achieved lower Mean Absolute Error (MAE) in 6 out of 8 blood groups.",
            "Both models correctly identified AB-Negative as mathematically unpredictable due to extreme sparsity.",
            "",
            "[Please insert your Comparison Table here in Google Slides]"
        ]
    )

    # --- Slide 7: Why Prophet ---
    add_bullet_slide(
        "The 'Glass-Box' Advantage in Healthcare",
        [
            "Healthcare requires trust; 'Black-box' models are often rejected by hospital administrators.",
            "Prophet natively decomposes forecasts into visible Trend, Weekly, and Yearly seasonalities.",
            "Allows staff to see exactly *why* a shortage is predicted.",
            "",
            "[Please insert your Prophet Forecast Graph here in Google Slides]"
        ]
    )

    # --- Slide 8: Snapshots ---
    add_bullet_slide(
        "Mid-Semester Operational Snapshots",
        [
            "Foundational React UI & Node.js authentication complete.",
            "FastAPI ML microservice successfully deployed locally.",
            "Returns JSON inventory alerts based on live parameter inputs.",
            "",
            "[Please insert your MERN UI and Swagger UI screenshots here in Google Slides]"
        ]
    )

    # --- Slide 9: Phase 2 Roadmap ---
    add_bullet_slide(
        "Phase 2: Enterprise Architecture Roadmap",
        [
            "Event-Driven Alerts: RabbitMQ for asynchronous emergency donor SMS/emails.",
            "Latency Optimization: Redis caching for sub-millisecond dashboard load times.",
            "Geospatial Routing: MongoDB GeoJSON to alert donors within a 10km radius.",
            "MLOps Pipeline: Automated CRON jobs for zero-downtime model retraining.",
            "DevOps: Containerizing the full stack via Docker."
        ]
    )

    # --- Slide 10: Conclusion ---
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Thank You"
    subtitle.text = "Open for Questions"

    # Save the presentation file
    output_filename = 'BTP_MidSem_Presentation.pptx'
    prs.save(output_filename)
    print(f"Success! Your presentation has been saved as '{output_filename}' in the current folder.")

if __name__ == "__main__":
    create_presentation()